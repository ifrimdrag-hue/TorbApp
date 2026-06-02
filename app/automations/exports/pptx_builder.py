"""PowerPoint pentru sedinta — 12 slide-uri corporate sobru.

Paleta: Midnight Executive (navy + light blue + gray)
Fonts: Cambria pentru titluri, Calibri pentru body.
Layout 16:9.
"""

from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from . import theme as T


# ─── Color helpers ───
def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


NAVY = rgb(T.NAVY)
LIGHT_BLUE = rgb(T.LIGHT_BLUE)
DARK_GRAY = rgb(T.DARK_GRAY)
MID_GRAY = rgb(T.MID_GRAY)
LIGHT_GRAY = rgb(T.LIGHT_GRAY)
WHITE = rgb(T.WHITE)


# ─── Slide dimensions (16:9) ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_text(tf, text, font=T.FONT_BODY, size=12, bold=False, italic=False,
              color=DARK_GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_text_box(slide, x, y, w, h, text, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    _set_text(box.text_frame, text, **kwargs)
    return box


def _add_filled_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_slide_title(slide, text):
    """Titlu standard de slide — sus, navy, aliniat stanga."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7))
    _set_text(box.text_frame, text, font=T.FONT_TITLE, size=28, bold=True, color=NAVY)
    return box


def _add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # 6 = blank in default templates
    return prs.slides.add_slide(blank_layout)


# ──────────────────────── SLIDE 1: COVER ────────────────────────
def _slide_cover(prs):
    slide = _add_blank_slide(prs)
    # Background navy on left third
    _add_filled_rect(slide, Inches(0), Inches(0), Inches(4.5), SLIDE_H, NAVY)
    # Vertical accent bar
    _add_filled_rect(slide, Inches(4.5), Inches(0), Inches(0.05), SLIDE_H, LIGHT_BLUE)

    # Title (right side, white space)
    _add_text_box(slide, Inches(5), Inches(2.4), Inches(7.8), Inches(1.4),
                  "Plan Campanii", font=T.FONT_TITLE, size=48, bold=True, color=NAVY)
    _add_text_box(slide, Inches(5), Inches(3.4), Inches(7.8), Inches(1.0),
                  "Mai 2026", font=T.FONT_TITLE, size=42, bold=False, color=LIGHT_BLUE)
    _add_text_box(slide, Inches(5), Inches(4.5), Inches(7.8), Inches(0.5),
                  "Marketing & E-commerce  ·  Propunere pentru sedinta",
                  size=14, italic=True, color=MID_GRAY)
    _add_text_box(slide, Inches(5), Inches(5.05), Inches(7.8), Inches(0.4),
                  date.today().strftime("%d %B %Y"), size=12, color=MID_GRAY)

    # On the navy band: small text
    _add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3.5), Inches(0.4),
                  "HUB AUTOMATIZARI", font=T.FONT_BODY, size=10, bold=True, color=WHITE)
    _add_text_box(slide, Inches(0.6), Inches(6.5), Inches(3.5), Inches(0.4),
                  "Brief executiv", size=11, italic=True, color=LIGHT_BLUE)


# ──────────────────────── SLIDE 2: CONTEXT ────────────────────────
def _slide_context(prs):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Context — De ce mai 2026?")

    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  "Trei momente comerciale distincte, complementare, cu potential demonstrat istoric.",
                  size=14, italic=True, color=MID_GRAY)

    # 3 cards
    moments = [
        ("Ziua Mamei", "duminica\n3 mai", "Gifting (peak comercial,\ncomparabil cu Martisor)", NAVY),
        ("Tranzitie\nprimavara→vara", "5 - 25 mai", "Lifestyle, refresh,\nterase deschise", LIGHT_BLUE),
        ("Pre-vara 'light'", "20 mai →\niunie", "Sanatate, fara zahar,\nhidratare", rgb("4A6FA5")),
    ]
    card_w = Inches(3.9)
    card_h = Inches(4.0)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    y = Inches(2.2)

    for i, (title, date_label, desc, color) in enumerate(moments):
        x = start_x + (card_w + gap) * i
        # Card background
        _add_filled_rect(slide, x, y, card_w, card_h, WHITE, line_color=rgb("D0D0D0"))
        # Top color band
        _add_filled_rect(slide, x, y, card_w, Inches(0.6), color)
        # Title (on color band)
        _add_text_box(slide, x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.4), Inches(0.5),
                      title.upper(), font=T.FONT_BODY, size=12, bold=True, color=WHITE)
        # Date big
        _add_text_box(slide, x + Inches(0.2), y + Inches(0.85), card_w - Inches(0.4), Inches(1.3),
                      date_label, font=T.FONT_TITLE, size=24, bold=True, color=NAVY)
        # Description
        _add_text_box(slide, x + Inches(0.2), y + Inches(2.4), card_w - Inches(0.4), Inches(1.4),
                      desc, size=12, color=DARK_GRAY)


# ──────────────────────── SLIDE 3: 3 CAMPANII OVERVIEW ────────────────────────
def _slide_overview(prs, campaigns):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Trei campanii complementare")

    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  "Fiecare adreseaza un moment comercial diferit. Fara canibalizare intre ele.",
                  size=14, italic=True, color=MID_GRAY)

    # 3 horizontal cards
    card_w = Inches(3.9)
    card_h = Inches(4.5)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    y = Inches(2.0)

    for i, c in enumerate(campaigns[:3]):
        x = start_x + (card_w + gap) * i
        _add_filled_rect(slide, x, y, card_w, card_h, WHITE, line_color=rgb("D0D0D0"))
        _add_filled_rect(slide, x, y, card_w, Inches(0.7), NAVY)

        _add_text_box(slide, x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.4), Inches(0.5),
                      f"  CAMPANIA {chr(65 + i)}", font=T.FONT_BODY, size=11, bold=True, color=LIGHT_BLUE)

        _add_text_box(slide, x + Inches(0.25), y + Inches(0.95), card_w - Inches(0.5), Inches(0.8),
                      c.get("name", ""), font=T.FONT_TITLE, size=15, bold=True, color=NAVY)

        # Period
        _add_text_box(slide, x + Inches(0.25), y + Inches(2.0), card_w - Inches(0.5), Inches(0.4),
                      f"📅 {c.get('date_start','')} → {c.get('date_end','')}",
                      size=11, color=DARK_GRAY)

        # Mechanic (truncated)
        mechanic = c.get("mechanic", "")
        if len(mechanic) > 100:
            mechanic = mechanic[:97] + "..."
        _add_text_box(slide, x + Inches(0.25), y + Inches(2.5), card_w - Inches(0.5), Inches(1.3),
                      f"🎯 {mechanic}", size=11, color=DARK_GRAY)

        # Budget big
        budget = c.get("budget_alloc") or 0
        _add_text_box(slide, x + Inches(0.25), y + Inches(3.7), card_w - Inches(0.5), Inches(0.7),
                      f"{int(budget)} RON",
                      font=T.FONT_TITLE, size=22, bold=True, color=LIGHT_BLUE)
        _add_text_box(slide, x + Inches(0.25), y + Inches(4.05), card_w - Inches(0.5), Inches(0.4),
                      "buget ads alocat", size=10, italic=True, color=MID_GRAY)


# ──────────────────────── SLIDE 4-6: DETAILED PER CAMPAIGN ────────────────────────
def _slide_campaign_detail(prs, c, idx):
    slide = _add_blank_slide(prs)
    type_labels = {"promo": "Promo", "gifting": "Gifting", "lansare": "Lansare", "sezonier": "Sezonier", "giveaway": "Giveaway"}
    channel_labels = {"shopify": "Shopify", "emag": "eMAG", "instagram": "Instagram", "facebook": "Facebook"}

    _add_slide_title(slide, f"Campania {chr(64 + idx)} — {c.get('name', '')}")

    # Subtitle line: type + period
    type_str = type_labels.get(c.get("type"), c.get("type", ""))
    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  f"{type_str}  ·  {c.get('date_start','')} → {c.get('date_end','')}",
                  size=14, italic=True, color=MID_GRAY)

    # Left column — Mecanica + Canale + Buget (stacked)
    left_x = Inches(0.6)
    right_x = Inches(7.0)
    col_w_l = Inches(6.0)
    col_w_r = Inches(5.7)
    y = Inches(2.0)

    # MECANICA box
    _add_filled_rect(slide, left_x, y, col_w_l, Inches(0.4), NAVY)
    _add_text_box(slide, left_x + Inches(0.15), y + Inches(0.05), col_w_l - Inches(0.3), Inches(0.3),
                  "MECANICA", font=T.FONT_BODY, size=10, bold=True, color=WHITE)
    _add_filled_rect(slide, left_x, y + Inches(0.4), col_w_l, Inches(1.4), LIGHT_GRAY,
                     line_color=rgb("D0D0D0"))
    _add_text_box(slide, left_x + Inches(0.2), y + Inches(0.5), col_w_l - Inches(0.4), Inches(1.2),
                  c.get("mechanic", ""), size=12, color=DARK_GRAY)

    # CANALE box (jos)
    y2 = y + Inches(2.0)
    _add_filled_rect(slide, left_x, y2, col_w_l, Inches(0.4), NAVY)
    _add_text_box(slide, left_x + Inches(0.15), y2 + Inches(0.05), col_w_l - Inches(0.3), Inches(0.3),
                  "CANALE", font=T.FONT_BODY, size=10, bold=True, color=WHITE)
    _add_filled_rect(slide, left_x, y2 + Inches(0.4), col_w_l, Inches(0.7), LIGHT_GRAY,
                     line_color=rgb("D0D0D0"))
    sales_ch = [ch for ch in c.get("channels", []) if ch in ("shopify", "emag")]
    post_ch = [ch for ch in c.get("channels", []) if ch in ("instagram", "facebook")]
    ch_str_parts = []
    if sales_ch:
        ch_str_parts.append("🛒 Vanzare: " + ", ".join(channel_labels.get(ch, ch) for ch in sales_ch))
    if post_ch:
        ch_str_parts.append("📣 Comunicare: " + ", ".join(channel_labels.get(ch, ch) for ch in post_ch))
    _add_text_box(slide, left_x + Inches(0.2), y2 + Inches(0.5), col_w_l - Inches(0.4), Inches(0.5),
                  "    ".join(ch_str_parts), size=11, color=DARK_GRAY)

    # BUGET box (jos jos)
    y3 = y2 + Inches(1.3)
    _add_filled_rect(slide, left_x, y3, col_w_l, Inches(0.4), NAVY)
    _add_text_box(slide, left_x + Inches(0.15), y3 + Inches(0.05), col_w_l - Inches(0.3), Inches(0.3),
                  "BUGET ADS", font=T.FONT_BODY, size=10, bold=True, color=WHITE)
    _add_filled_rect(slide, left_x, y3 + Inches(0.4), col_w_l, Inches(0.9), WHITE,
                     line_color=rgb("D0D0D0"))
    budget = c.get("budget_alloc") or 0
    _add_text_box(slide, left_x + Inches(0.3), y3 + Inches(0.55), col_w_l - Inches(0.6), Inches(0.7),
                  f"{int(budget)} RON",
                  font=T.FONT_TITLE, size=28, bold=True, color=LIGHT_BLUE)

    # ─── Right column — Produse incluse ───
    _add_filled_rect(slide, right_x, y, col_w_r, Inches(0.4), NAVY)
    _add_text_box(slide, right_x + Inches(0.15), y + Inches(0.05), col_w_r - Inches(0.3), Inches(0.3),
                  f"PRODUSE INCLUSE ({len(c.get('products', []))})",
                  font=T.FONT_BODY, size=10, bold=True, color=WHITE)
    prod_h = Inches(4.2)
    _add_filled_rect(slide, right_x, y + Inches(0.4), col_w_r, prod_h, WHITE,
                     line_color=rgb("D0D0D0"))

    prods = (c.get("products") or [])[:8]
    py = y + Inches(0.55)
    for p in prods:
        name = p.get("name", p.get("sku", ""))
        if len(name) > 55:
            name = name[:52] + "..."
        _add_text_box(slide, right_x + Inches(0.2), py, col_w_r - Inches(0.4), Inches(0.32),
                      f"•  {name}", size=10, color=DARK_GRAY)
        if p.get("qty_needed"):
            _add_text_box(slide, right_x + Inches(0.2), py + Inches(0.28),
                          col_w_r - Inches(0.4), Inches(0.22),
                          f"     necesar: {p['qty_needed']} buc",
                          size=9, italic=True, color=MID_GRAY)
            py += Inches(0.55)
        else:
            py += Inches(0.4)
    if len(c.get("products", [])) > 8:
        _add_text_box(slide, right_x + Inches(0.2), py + Inches(0.1),
                      col_w_r - Inches(0.4), Inches(0.3),
                      f"   ... si inca {len(c['products']) - 8} produse",
                      size=10, italic=True, color=MID_GRAY)


# ──────────────────────── SLIDE 7: BUGET TOTAL CU CHART ────────────────────────
def _slide_budget(prs, campaigns):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Buget propus")

    total = sum(c.get("budget_alloc") or 0 for c in campaigns)
    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  f"Total ads: {int(total)} RON  ·  ~200 EUR  ·  Distributie pe campanii:",
                  size=14, italic=True, color=MID_GRAY)

    # Pie chart left
    chart_data = CategoryChartData()
    chart_data.categories = [c.get("name", f"Campanie {i}") for i, c in enumerate(campaigns, 1)]
    chart_data.add_series("Buget alocat", [c.get("budget_alloc") or 0 for c in campaigns])

    chart_x, chart_y = Inches(0.6), Inches(2.0)
    chart_w, chart_h = Inches(6.0), Inches(4.8)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, chart_x, chart_y, chart_w, chart_h, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.show_percentage = False
    dl.font.size = Pt(11)
    dl.font.bold = True

    # Text right column
    rx = Inches(7.0)
    _add_text_box(slide, rx, Inches(2.0), Inches(5.7), Inches(0.5),
                  "BUGET ADS PE CAMPANIE", font=T.FONT_BODY, size=11, bold=True, color=NAVY)

    ry = Inches(2.6)
    for c in campaigns:
        name = c.get("name", "")
        if len(name) > 36:
            name = name[:33] + "..."
        budget = int(c.get("budget_alloc") or 0)
        _add_text_box(slide, rx, ry, Inches(4.5), Inches(0.4),
                      name, size=12, color=DARK_GRAY)
        _add_text_box(slide, rx + Inches(4.5), ry, Inches(1.2), Inches(0.4),
                      f"{budget} RON", size=12, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)
        ry += Inches(0.5)

    # Total
    ry += Inches(0.2)
    _add_filled_rect(slide, rx, ry, Inches(5.7), Inches(0.05), NAVY)
    ry += Inches(0.15)
    _add_text_box(slide, rx, ry, Inches(4.5), Inches(0.5),
                  "TOTAL ads", font=T.FONT_BODY, size=14, bold=True, color=NAVY)
    _add_text_box(slide, rx + Inches(4.5), ry, Inches(1.2), Inches(0.5),
                  f"{int(total)} RON", font=T.FONT_TITLE, size=18, bold=True, color=NAVY,
                  align=PP_ALIGN.RIGHT)

    ry += Inches(0.7)
    _add_text_box(slide, rx, ry, Inches(5.7), Inches(0.7),
                  "Buget producere creativa\n(foto, video, design): 1.500-2.500 RON separat — "
                  "asseturi reutilizabile peste 3 campanii.",
                  size=10, italic=True, color=MID_GRAY)


# ──────────────────────── SLIDE 8: TIMELINE / GANTT ────────────────────────
def _slide_timeline(prs, campaigns):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Calendar campanii")
    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  "Ritm crescendo: A intens si scurt → B constant lung → C educational long-tail.",
                  size=14, italic=True, color=MID_GRAY)

    # Timeline axis
    tl_x = Inches(2.5)
    tl_w = Inches(10.0)
    tl_y = Inches(2.5)
    # Date range: 28 apr → 12 iun (45 zile)
    from datetime import datetime as dt
    start = dt(2026, 4, 28)
    end = dt(2026, 6, 12)
    total_days = (end - start).days

    def day_to_x(d_str):
        d = dt.strptime(d_str, "%Y-%m-%d")
        offset = (d - start).days
        return tl_x + Emu(int((tl_w / total_days) * offset))

    # Month labels at top
    months = [("APRILIE", dt(2026, 4, 28)), ("MAI", dt(2026, 5, 1)), ("IUNIE", dt(2026, 6, 1))]
    for label, d in months:
        x = tl_x + Emu(int((tl_w / total_days) * (d - start).days))
        _add_text_box(slide, x, tl_y - Inches(0.5), Inches(2.0), Inches(0.4),
                      label, font=T.FONT_BODY, size=10, bold=True, color=MID_GRAY)

    # Axis line
    _add_filled_rect(slide, tl_x, tl_y, tl_w, Inches(0.02), MID_GRAY)

    # Date markers (every week)
    week_starts = [dt(2026, 4, 28), dt(2026, 5, 5), dt(2026, 5, 12), dt(2026, 5, 19),
                   dt(2026, 5, 26), dt(2026, 6, 2), dt(2026, 6, 9)]
    for d in week_starts:
        x = tl_x + Emu(int((tl_w / total_days) * (d - start).days))
        _add_filled_rect(slide, x, tl_y - Inches(0.05), Inches(0.02), Inches(0.15), MID_GRAY)
        _add_text_box(slide, x - Inches(0.4), tl_y + Inches(0.1), Inches(0.8), Inches(0.3),
                      d.strftime("%d.%m"), size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Campaign bars
    bar_h = Inches(0.55)
    bar_y = tl_y + Inches(0.7)
    colors = [NAVY, LIGHT_BLUE, rgb("4A6FA5")]
    for i, c in enumerate(campaigns[:3]):
        x_start = day_to_x(c["date_start"])
        x_end = day_to_x(c["date_end"])
        bar_w = x_end - x_start

        y_offset = bar_y + Inches(0.85) * i

        # Campaign label (left)
        _add_text_box(slide, Inches(0.6), y_offset + Inches(0.08), Inches(1.7), Inches(0.4),
                      c.get("name", "").split("—")[0].strip(),
                      size=11, bold=True, color=DARK_GRAY)

        # Bar
        _add_filled_rect(slide, x_start, y_offset, bar_w, bar_h, colors[i])

        # Bar label inside (or beside if too narrow)
        # Check if bar is wide enough for label
        if bar_w > Inches(2.5):
            mechanic = c.get("mechanic", "")
            if len(mechanic) > 35:
                mechanic = mechanic[:32] + "..."
            _add_text_box(slide, x_start + Inches(0.15), y_offset + Inches(0.12),
                          bar_w - Inches(0.3), Inches(0.3),
                          mechanic, size=9, color=WHITE)

        # Date range under bar
        _add_text_box(slide, x_start, y_offset + bar_h + Inches(0.02),
                      Inches(3.0), Inches(0.3),
                      f"{c['date_start'][8:]}.{c['date_start'][5:7]} → {c['date_end'][8:]}.{c['date_end'][5:7]}",
                      size=8, italic=True, color=MID_GRAY)


# ──────────────────────── SLIDE 9: KPI TINTE ────────────────────────
def _slide_kpis(prs):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Rezultate asteptate")
    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  "Tinte agregate pentru luna mai 2026:",
                  size=14, italic=True, color=MID_GRAY)

    # 6 stat cards, 3x2 grid
    stats = [
        ("350-450", "comenzi B2C totale", LIGHT_BLUE),
        ("25-40K", "RON revenue B2C", NAVY),
        ("+15-25%", "AOV vs aprilie", LIGHT_BLUE),
        (">28%", "contribution margin", NAVY),
        ("+200-400", "abonati newsletter", LIGHT_BLUE),
        ("60-70%", "lichidare stoc Torras", NAVY),
    ]

    card_w = Inches(3.9)
    card_h = Inches(2.4)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    start_y = Inches(2.0)

    for idx, (val, label, color) in enumerate(stats):
        col = idx % 3
        row = idx // 3
        x = start_x + (card_w + gap) * col
        y = start_y + (card_h + gap) * row

        _add_filled_rect(slide, x, y, card_w, card_h, WHITE, line_color=rgb("D0D0D0"))
        _add_filled_rect(slide, x, y, Inches(0.15), card_h, color)

        _add_text_box(slide, x + Inches(0.4), y + Inches(0.4), card_w - Inches(0.6), Inches(1.2),
                      val, font=T.FONT_TITLE, size=44, bold=True, color=color)
        _add_text_box(slide, x + Inches(0.4), y + Inches(1.55), card_w - Inches(0.6), Inches(0.7),
                      label, size=13, color=DARK_GRAY)


# ──────────────────────── SLIDE 10: RISCURI ────────────────────────
def _slide_risks(prs):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Riscuri & mitigari")

    risks = [
        ("Stoc insuficient Tea Book de Ziua Mamei", "Medie",
         "Verific stoc prin Hub inainte de lansare. Limitez bundle la N buc disponibile."),
        ("Canibalizare intre B si C", "Mica",
         "Audiente diferite — B = lifestyle general, C = health-conscious."),
        ("Vreme rece prelungita", "Mica-Medie",
         "Mesaj poate fi pivotat la 'indoor coziness' / ceaiuri calde."),
        ("Foto-sesiunea intarziata", "Medie",
         "Programare imediat dupa aprobare. Buffer 2 zile inainte de lansare."),
        ("App Shopify pentru 2+1 nu functioneaza", "Identificat",
         "Plan B activ: cod discount manual la checkout."),
    ]

    # Header row
    y = Inches(1.5)
    headers_w = [Inches(5.5), Inches(1.8), Inches(5.4)]
    headers_x = [Inches(0.6), Inches(0.6) + headers_w[0], Inches(0.6) + headers_w[0] + headers_w[1]]
    for i, h in enumerate(["RISC", "PROBABILITATE", "MITIGARE"]):
        _add_filled_rect(slide, headers_x[i], y, headers_w[i], Inches(0.5), NAVY)
        _add_text_box(slide, headers_x[i] + Inches(0.2), y + Inches(0.1), headers_w[i] - Inches(0.4), Inches(0.3),
                      h, font=T.FONT_BODY, size=11, bold=True, color=WHITE)

    y += Inches(0.5)
    row_h = Inches(0.95)
    for idx, (risk, prob, mit) in enumerate(risks):
        bg = LIGHT_GRAY if idx % 2 == 0 else WHITE
        for i in range(3):
            _add_filled_rect(slide, headers_x[i], y, headers_w[i], row_h, bg, line_color=rgb("D0D0D0"))

        _add_text_box(slide, headers_x[0] + Inches(0.2), y + Inches(0.15),
                      headers_w[0] - Inches(0.4), row_h - Inches(0.3),
                      risk, size=11, color=DARK_GRAY)
        _add_text_box(slide, headers_x[1] + Inches(0.2), y + Inches(0.15),
                      headers_w[1] - Inches(0.4), row_h - Inches(0.3),
                      prob, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_text_box(slide, headers_x[2] + Inches(0.2), y + Inches(0.15),
                      headers_w[2] - Inches(0.4), row_h - Inches(0.3),
                      mit, size=10, color=DARK_GRAY)
        y += row_h


# ──────────────────────── SLIDE 11: ASKS ────────────────────────
def _slide_asks(prs):
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, "Decizii necesare astazi")
    _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                  "Pentru lansare conform calendarului, am nevoie de aprobare pe urmatoarele:",
                  size=14, italic=True, color=MID_GRAY)

    asks = [
        ("Aprobare buget total ads",
         "1.000 RON pentru toate cele 3 campanii (200 EUR)"),
        ("Aprobare buget producere creativa",
         "1.500-2.500 RON foto+video, sau folosim asseturile existente"),
        ("Confirmare stoc Basilur Tea Book",
         "Vol I/II/III suficient pentru bundle-ul de Ziua Mamei (~80-100 buc)"),
        ("Coordonare cu KAM",
         "Campania C (Verile Usor) e relevanta si pentru clienti B2B (gym, cafenele)?"),
        ("Validare ROI istoric",
         "Date Ziua Mamei 2025 pentru calibrare asteptari realiste"),
    ]

    y = Inches(1.9)
    for i, (title, detail) in enumerate(asks):
        # Number circle
        circle_size = Inches(0.7)
        _add_filled_rect(slide, Inches(0.6), y, circle_size, circle_size, NAVY)
        _add_text_box(slide, Inches(0.6), y + Inches(0.13), circle_size, Inches(0.5),
                      str(i + 1), font=T.FONT_TITLE, size=22, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
        # Text
        _add_text_box(slide, Inches(1.6), y, Inches(11.0), Inches(0.4),
                      title, font=T.FONT_BODY, size=14, bold=True, color=NAVY)
        _add_text_box(slide, Inches(1.6), y + Inches(0.4), Inches(11.0), Inches(0.4),
                      detail, size=11, color=DARK_GRAY)
        y += Inches(1.05)


# ──────────────────────── SLIDE 12: NEXT STEPS / CLOSE ────────────────────────
def _slide_close(prs):
    slide = _add_blank_slide(prs)
    # Background navy
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Vertical accent
    _add_filled_rect(slide, Inches(0.6), Inches(1.5), Inches(0.05), Inches(4.5), LIGHT_BLUE)

    _add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.5), Inches(0.7),
                  "Next steps imediate", font=T.FONT_TITLE, size=32, bold=True, color=WHITE)

    _add_text_box(slide, Inches(1.0), Inches(2.3), Inches(11.5), Inches(0.4),
                  "Dupa aprobarea sedintei", size=14, italic=True, color=LIGHT_BLUE)

    steps = [
        ("30 aprilie", "Lansare Mama Mea + listare bundle Shopify"),
        ("1 mai", "Foto-sesiune programata + setup ads Meta (~18:00)"),
        ("2 mai", "Email blast 'last chance' Mother's Day"),
        ("3 mai", "Mother's Day — monitoring intensiv"),
        ("5 mai", "Lansare campania B (Sezonul Schimba)"),
        ("Saptamanal", "Check-in cu directorul comercial — vineri 15:00"),
    ]

    y = Inches(3.2)
    for date_str, action in steps:
        _add_text_box(slide, Inches(1.0), y, Inches(2.0), Inches(0.4),
                      date_str, font=T.FONT_BODY, size=13, bold=True, color=LIGHT_BLUE)
        _add_text_box(slide, Inches(3.0), y, Inches(9.5), Inches(0.4),
                      action, size=13, color=WHITE)
        y += Inches(0.55)

    # Footer
    _add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
                  "Plan generat din Hub Automatizari — toate datele sunt live in sistem.",
                  size=10, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ──────────────────────── BUILD ────────────────────────
def build_pptx(campaigns: list[dict]) -> bytes:
    """Construieste prezentarea PPTX completa din lista de campanii."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Sortam campaniile dupa data de start pentru consistenta
    campaigns = sorted(campaigns, key=lambda c: c.get("date_start", ""))

    _slide_cover(prs)
    _slide_context(prs)
    _slide_overview(prs, campaigns)
    for i, c in enumerate(campaigns[:3], start=1):
        _slide_campaign_detail(prs, c, i)
    _slide_budget(prs, campaigns)
    _slide_timeline(prs, campaigns)
    _slide_kpis(prs)
    _slide_risks(prs)
    _slide_asks(prs)
    _slide_close(prs)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
