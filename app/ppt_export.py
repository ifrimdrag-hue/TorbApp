"""
PPT export pentru Torb Logistic.
Generează prezentări PowerPoint cu date comerciale pentru board.
"""

import io
import datetime
from flask import send_file

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    PPTX_OK = True

    # ── Culori Torb ──────────────────────────────────────────────────────────
    C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
    C_BLUE   = RGBColor(0x16, 0x21, 0x3e)
    C_ACCENT = RGBColor(0x0f, 0x3d, 0x5f)
    C_GREEN  = RGBColor(0x2e, 0xcc, 0x71)
    C_RED    = RGBColor(0xe7, 0x4c, 0x3c)
    C_ORANGE = RGBColor(0xf3, 0x9c, 0x12)
    C_WHITE  = RGBColor(0xff, 0xff, 0xff)
    C_LGRAY  = RGBColor(0xf5, 0xf6, 0xfa)
    C_MGRAY  = RGBColor(0xcc, 0xcc, 0xcc)
    C_TEXT   = RGBColor(0x2d, 0x35, 0x36)

except ImportError:
    PPTX_OK = False
    C_DARK = C_BLUE = C_ACCENT = C_GREEN = C_RED = C_ORANGE = None
    C_WHITE = C_LGRAY = C_MGRAY = C_TEXT = None


def _check():
    if not PPTX_OK:
        raise RuntimeError("python-pptx nu este instalat. Rulează: pip install python-pptx")


def _slide_w():
    return Inches(13.33)


def _slide_h():
    return Inches(7.5)


def _prs():
    """New widescreen (16:9) presentation."""
    prs = Presentation()
    prs.slide_width  = _slide_w()
    prs.slide_height = _slide_h()
    return prs


def _blank(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    _fill(shape, color)
    shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              font_size=12, bold=False, color=None, align=None,
              wrap=True):
    if color is None:
        color = C_TEXT
    if align is None:
        align = PP_ALIGN.LEFT
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _header_bar(slide, title, subtitle=""):
    """Dark header bar across the top."""
    _add_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    _add_text(slide, title, 0.3, 0.1, 10, 0.55,
              font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_text(slide, subtitle, 0.3, 0.65, 10, 0.4,
                  font_size=11, color=C_MGRAY, align=PP_ALIGN.LEFT)


def _footer(slide, text="Torb Logistic — Confidențial"):
    _add_rect(slide, 0, 7.2, 13.33, 0.3, C_BLUE)
    _add_text(slide, text, 0.3, 7.22, 12, 0.25,
              font_size=8, color=C_MGRAY)


def _kpi_card(slide, left, top, width, label, value, delta=None, delta_positive=None):
    """Single KPI card box."""
    _add_rect(slide, left, top, width, 1.5, C_LGRAY)
    # accent bar left
    _add_rect(slide, left, top, 0.06, 1.5, C_ACCENT)
    _add_text(slide, label, left + 0.12, top + 0.08, width - 0.2, 0.3,
              font_size=9, color=C_ACCENT, bold=True)
    _add_text(slide, value, left + 0.12, top + 0.35, width - 0.2, 0.65,
              font_size=18, bold=True, color=C_TEXT)
    if delta is not None:
        col = C_GREEN if delta_positive else C_RED
        arrow = "▲" if delta_positive else "▼"
        _add_text(slide, f"{arrow} {delta}", left + 0.12, top + 1.05, width - 0.2, 0.35,
                  font_size=10, bold=True, color=col)


def _fmt_ron(v):
    if v is None:
        return "—"
    try:
        v = float(v)
        if abs(v) >= 1_000_000:
            return f"{v/1_000_000:.2f}M RON"
        return f"{int(v):,} RON".replace(",", ".")
    except Exception:
        return str(v)


def _fmt_pct(v, suffix="%"):
    if v is None:
        return "—"
    try:
        return f"{float(v):.1f}{suffix}"
    except Exception:
        return str(v)


def _add_table(slide, rows, headers, left, top, width, height, font_size=8):
    """Add a data table to a slide."""
    if not rows:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols,
        Inches(left), Inches(top), Inches(width), Inches(height)).table
    tbl.first_row = True

    col_w = width / n_cols
    for i in range(n_cols):
        tbl.columns[i].width = Inches(col_w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    # Data rows
    for i, row in enumerate(rows):
        for j, key in enumerate(row):
            cell = tbl.cell(i + 1, j)
            val = row[key] if isinstance(row, dict) else row[j]
            cell.text = str(val) if val is not None else "—"
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_LGRAY if i % 2 == 0 else C_WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = C_TEXT

    return tbl


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(prs, title, subtitle, an):
    slide = _blank(prs)
    # Full dark background
    _add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    _add_rect(slide, 0, 0, 0.25, 7.5, C_GREEN)
    _add_text(slide, "TORB LOGISTIC", 0.5, 1.5, 12, 0.8,
              font_size=14, bold=True, color=C_MGRAY)
    _add_text(slide, title, 0.5, 2.3, 12, 1.5,
              font_size=36, bold=True, color=C_WHITE)
    _add_text(slide, subtitle, 0.5, 3.9, 12, 0.6,
              font_size=16, color=C_MGRAY)
    _add_text(slide, str(an), 0.5, 4.6, 4, 0.5,
              font_size=14, bold=True, color=C_GREEN)
    today = datetime.date.today().strftime("%d.%m.%Y")
    _add_text(slide, f"Generat: {today}", 0.5, 6.8, 5, 0.4,
              font_size=9, color=C_MGRAY)
    return slide


def _slide_kpi_overview(prs, an, cy, py, delta_vn, delta_mb, delta_mn, delta_mpct):
    slide = _blank(prs)
    _header_bar(slide, f"Performanță comercială {an}", "KPI Cheie — comparativ an anterior")
    _footer(slide)

    cards = [
        ("Valoare Factură Netă", _fmt_ron(cy.get('val_neta')),
         _fmt_pct(delta_vn, "%"), (delta_vn or 0) >= 0),
        ("Marjă Brută", _fmt_ron(cy.get('marja_bruta')),
         _fmt_pct(delta_mb, "%"), (delta_mb or 0) >= 0),
        ("Marjă Netă", _fmt_ron(cy.get('marja_neta')),
         _fmt_pct(delta_mn, "%"), (delta_mn or 0) >= 0),
        ("Marjă Netă %", _fmt_pct(cy.get('marja_pct')),
         f"{delta_mpct:+.1f}pp" if delta_mpct is not None else "—",
         (delta_mpct or 0) >= 0),
    ]
    x_start = 0.3
    card_w = 3.0
    gap = 0.22
    for i, (label, value, delta, pos) in enumerate(cards):
        _kpi_card(slide, x_start + i * (card_w + gap), 1.25, card_w,
                  label, value, delta, pos)

    # Active clients
    n_clients = cy.get('clienti_activi', 0)
    n_py = py.get('clienti_activi', 0) if py else 0
    delta_c = n_clients - n_py
    _add_text(slide, "Clienți Activi", 0.3, 3.0, 6, 0.35,
              font_size=11, bold=True, color=C_ACCENT)
    _add_text(slide, str(n_clients), 0.3, 3.35, 3, 0.7,
              font_size=32, bold=True, color=C_TEXT)
    col = C_GREEN if delta_c >= 0 else C_RED
    arrow = "▲" if delta_c >= 0 else "▼"
    _add_text(slide, f"{arrow} {abs(delta_c)} față de {an-1}", 0.3, 4.0, 5, 0.35,
              font_size=11, bold=True, color=col)

    # MB% vs MN% mini comparison
    mb_pct = cy.get('marja_pct') or 0
    mn_pct = (cy.get('marja_neta') or 0) / (cy.get('val_neta') or 1) * 100
    _add_text(slide, "Marjă Brută %", 7.0, 3.0, 5, 0.3, font_size=10, bold=True, color=C_ACCENT)
    _add_text(slide, _fmt_pct(mb_pct), 7.0, 3.3, 3, 0.5, font_size=28, bold=True, color=C_TEXT)
    _add_text(slide, "Marjă Netă %", 10.2, 3.0, 3, 0.3, font_size=10, bold=True, color=C_ACCENT)
    _add_text(slide, _fmt_pct(mn_pct), 10.2, 3.3, 3, 0.5, font_size=28, bold=True, color=C_RED if mn_pct < 5 else C_TEXT)

    return slide


def _slide_agents_table(prs, an, agents_data):
    slide = _blank(prs)
    _header_bar(slide, f"Performanță Agenți {an}", "Valoare Netă / Marjă Brută / Marjă Netă")
    _footer(slide)

    headers = ["Agent", "Val. Netă", "MB RON", "MB %", "MN RON", "MN %", "Clienți", "vs AN-1 MN"]
    rows = []
    for a in agents_data[:15]:
        delta_str = ""
        delta_mn = a.get('delta_mn_pct')
        if delta_mn is not None:
            arrow = "▲" if delta_mn >= 0 else "▼"
            delta_str = f"{arrow}{abs(delta_mn):.1f}%"
        rows.append({
            "Agent": (a.get('agent') or "—")[:25],
            "Val. Netă": _fmt_ron(a.get('val_neta')),
            "MB RON": _fmt_ron(a.get('marja_bruta')),
            "MB %": _fmt_pct(a.get('marja_bruta_pct')),
            "MN RON": _fmt_ron(a.get('marja_neta')),
            "MN %": _fmt_pct(a.get('marja_neta_pct')),
            "Clienți": str(a.get('nr_clienti') or 0),
            "vs AN-1 MN": delta_str,
        })

    n = min(len(rows), 15)
    tbl_h = min(5.6, 0.38 * (n + 1))
    _add_table(slide, rows, list(rows[0].keys()) if rows else headers,
               0.3, 1.2, 12.7, tbl_h, font_size=9)
    return slide


def _slide_top_clients(prs, an, clients_data, title="Top Clienți după Marjă Netă"):
    slide = _blank(prs)
    _header_bar(slide, title, str(an))
    _footer(slide)

    headers = ["#", "Client", "Agent", "Val. Netă", "MB %", "MN RON", "MN %"]
    rows = []
    for i, c in enumerate(clients_data[:15], 1):
        rows.append({
            "#": str(i),
            "Client": (c.get('client') or "—")[:30],
            "Agent": (c.get('agent') or "—")[:15],
            "Val. Netă": _fmt_ron(c.get('val_neta')),
            "MB %": _fmt_pct(c.get('marja_bruta_pct')),
            "MN RON": _fmt_ron(c.get('marja_neta')),
            "MN %": _fmt_pct(c.get('marja_neta_pct')),
        })

    n = min(len(rows), 15)
    tbl_h = min(5.6, 0.38 * (n + 1))
    _add_table(slide, rows, list(rows[0].keys()) if rows else headers,
               0.3, 1.2, 12.7, tbl_h, font_size=9)
    return slide


def _slide_agent_detail(prs, agent_name, an, kpi, clients_data, brands_data):
    slide = _blank(prs)
    _header_bar(slide, f"Agent: {agent_name}", f"Performanță detaliată {an}")
    _footer(slide)

    # KPI mini cards
    cards = [
        ("Val. Netă", _fmt_ron(kpi.get('val_neta'))),
        ("Marjă Brută", f"{_fmt_ron(kpi.get('marja_bruta'))} / {_fmt_pct(kpi.get('marja_pct'))}"),
        ("Marjă Netă", f"{_fmt_ron(kpi.get('marja_neta'))} / {_fmt_pct(kpi.get('marja_neta_pct') if kpi.get('marja_neta_pct') else None)}"),
        ("Clienți Activi", str(kpi.get('clienti_activi') or 0)),
    ]
    for i, (lbl, val) in enumerate(cards):
        x = 0.3 + i * 3.15
        _add_rect(slide, x, 1.2, 3.0, 0.8, C_LGRAY)
        _add_text(slide, lbl, x + 0.1, 1.25, 2.8, 0.28, font_size=8, color=C_ACCENT, bold=True)
        _add_text(slide, val, x + 0.1, 1.5, 2.8, 0.45, font_size=13, bold=True, color=C_TEXT)

    # Clients table (left)
    _add_text(slide, "Clienți", 0.3, 2.15, 7, 0.3, font_size=10, bold=True, color=C_DARK)
    c_headers = ["Client", "Val. Netă", "MB%", "MN RON", "MN%"]
    c_rows = [{
        "Client": (c.get('client') or "—")[:22],
        "Val. Netă": _fmt_ron(c.get('val_neta')),
        "MB%": _fmt_pct(c.get('marja_bruta_pct')),
        "MN RON": _fmt_ron(c.get('marja_neta')),
        "MN%": _fmt_pct(c.get('marja_neta_pct')),
    } for c in (clients_data or [])[:8]]
    if c_rows:
        _add_table(slide, c_rows, c_headers, 0.3, 2.45, 6.9, min(4.5, 0.4*(len(c_rows)+1)), font_size=8)

    # Brands table (right)
    _add_text(slide, "Brand Mix", 7.4, 2.15, 5.6, 0.3, font_size=10, bold=True, color=C_DARK)
    b_headers = ["Brand", "Val. Netă", "MB%", "MN%"]
    b_rows = [{
        "Brand": b.get('furnizor') or "—",
        "Val. Netă": _fmt_ron(b.get('val_neta')),
        "MB%": _fmt_pct(b.get('marja_bruta_pct')),
        "MN%": _fmt_pct(b.get('marja_neta_pct')),
    } for b in (brands_data or [])[:8]]
    if b_rows:
        _add_table(slide, b_rows, b_headers, 7.4, 2.45, 5.6, min(4.5, 0.4*(len(b_rows)+1)), font_size=8)

    return slide


def _slide_client_detail(prs, client_name, an, kpi, products_data, yearly_data):
    slide = _blank(prs)
    _header_bar(slide, f"Client: {client_name[:40]}", f"Profil de profitabilitate {an}")
    _footer(slide)

    # KPI
    cards = [
        ("Val. Netă", _fmt_ron(kpi.get('val_neta_total') or kpi.get('val_neta'))),
        ("Marjă Brută", _fmt_pct(kpi.get('marja_pct'))),
        ("Marjă Netă", _fmt_ron(kpi.get('marja_neta'))),
        ("Nr. Facturi", str(kpi.get('nr_facturi') or 0)),
    ]
    for i, (lbl, val) in enumerate(cards):
        x = 0.3 + i * 3.15
        _add_rect(slide, x, 1.2, 3.0, 0.8, C_LGRAY)
        _add_text(slide, lbl, x + 0.1, 1.25, 2.8, 0.28, font_size=8, color=C_ACCENT, bold=True)
        _add_text(slide, val, x + 0.1, 1.5, 2.8, 0.45, font_size=13, bold=True, color=C_TEXT)

    # Products table (left)
    _add_text(slide, f"Produse cumpărate {an}", 0.3, 2.15, 7, 0.3, font_size=10, bold=True, color=C_DARK)
    p_headers = ["Produs", "Brand", "VN", "MB%", "MN%"]
    p_rows = [{
        "Produs": (r.get('sku') or "—")[:28],
        "Brand": r.get('furnizor') or "—",
        "VN": _fmt_ron(r.get('val_neta')),
        "MB%": _fmt_pct(r.get('marja_bruta_pct')),
        "MN%": _fmt_pct(r.get('marja_neta_pct')),
    } for r in (products_data or [])[:8]]
    if p_rows:
        _add_table(slide, p_rows, p_headers, 0.3, 2.45, 7.6, min(4.5, 0.4*(len(p_rows)+1)), font_size=8)

    # Yearly table (right)
    _add_text(slide, "Evoluție anuală", 8.1, 2.15, 4.9, 0.3, font_size=10, bold=True, color=C_DARK)
    y_headers = ["An", "Val. Netă", "MB%", "MN%"]
    y_rows = [{
        "An": str(r.get('an') or ""),
        "Val. Netă": _fmt_ron(r.get('val_neta')),
        "MB%": _fmt_pct(r.get('marja_bruta_pct')),
        "MN%": _fmt_pct(r.get('marja_neta_pct')),
    } for r in (yearly_data or [])]
    if y_rows:
        _add_table(slide, y_rows, y_headers, 8.1, 2.45, 4.9, min(4.5, 0.4*(len(y_rows)+1)), font_size=9)

    return slide


def _slide_risk(prs, an, kaufland_pct, bogdan_pct, churn_list):
    slide = _blank(prs)
    _header_bar(slide, "Riscuri Comerciale", "Concentrare și Churn")
    _footer(slide)

    # Concentration risk
    _add_text(slide, "⚠ Risc Concentrare", 0.3, 1.2, 6, 0.4, font_size=14, bold=True, color=C_ORANGE)

    risk_items = [
        ("Kaufland", kaufland_pct, "din total vânzări", 42.0),
        ("Bogdan Dragnea", bogdan_pct, "din total vânzări", 50.0),
    ]
    for i, (name, pct, label, threshold) in enumerate(risk_items):
        y = 1.75 + i * 1.5
        _add_rect(slide, 0.3, y, 6, 1.2, C_LGRAY)
        _add_text(slide, name, 0.45, y + 0.1, 5.5, 0.35, font_size=12, bold=True, color=C_TEXT)
        col = C_RED if (pct or 0) > threshold else C_ORANGE
        _add_text(slide, _fmt_pct(pct), 0.45, y + 0.4, 3, 0.55, font_size=26, bold=True, color=col)
        _add_text(slide, label, 0.45, y + 0.9, 5, 0.3, font_size=9, color=C_MGRAY)

    # Churn
    _add_text(slide, f"Clienți Inactivi (>60 zile): {len(churn_list)}",
              7.0, 1.2, 6, 0.4, font_size=14, bold=True, color=C_RED)
    ch_headers = ["Client", "Zile Inactiv", "Ultima Comandă"]
    ch_rows = [{
        "Client": (c.get('client') or "—")[:28],
        "Zile Inactiv": str(c.get('zile_inactiv') or "—"),
        "Ultima Comandă": c.get('ultima_comanda') or "—",
    } for c in (churn_list or [])[:10]]
    if ch_rows:
        _add_table(slide, ch_rows, ch_headers, 7.0, 1.7, 6.0, min(5.0, 0.4*(len(ch_rows)+1)), font_size=9)
    return slide


MONTHS_RO = ['Ian', 'Feb', 'Mar', 'Apr', 'Mai', 'Iun',
             'Iul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']


def _slide_chart_trend(prs, an, trend_by_year):
    """Bar chart cu trend lunar pe 3 ani (date YTD din monthly_trend)."""
    slide = _blank(prs)
    _header_bar(slide, f"Trend Lunar Vânzări — Comparativ {an-2}/{an-1}/{an}",
                "Val Netă RON pe lună")
    _footer(slide)

    data = CategoryChartData()
    data.categories = MONTHS_RO
    for yr in sorted(trend_by_year.keys()):
        data.add_series(str(yr), trend_by_year[yr])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(1.1),
        Inches(12.7), Inches(5.8),
        data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    return slide


def _slide_chart_brands(prs, an, brands_data):
    """Donut chart cu mix branduri pentru `an`."""
    slide = _blank(prs)
    _header_bar(slide, f"Mix Branduri {an}", "Distribuție Val Netă pe furnizor")
    _footer(slide)

    items = [(b.get('furnizor') or '—', b.get('val_neta') or 0)
             for b in brands_data or []]
    items = [(label, val) for label, val in items if val > 0]

    if items:
        data = CategoryChartData()
        data.categories = [b[0] for b in items]
        data.add_series('Val Netă', [b[1] for b in items])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.1),
            Inches(8.5), Inches(5.8),
            data,
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    # Tabel sumar lateral
    total = sum(v for _, v in items) or 1
    rows = [{
        "Brand":   label,
        "Val Netă": _fmt_ron(val),
        "%":       f"{val * 100.0 / total:.1f}%",
    } for label, val in items[:12]]
    if rows:
        _add_table(slide, rows, list(rows[0].keys()),
                   9.3, 1.4, 3.8, min(5.5, 0.35*(len(rows)+1)), font_size=9)
    return slide


def _slide_chart_channels(prs, an, channels_data):
    """Horizontal bar chart cu vânzări pe canale/agenți."""
    slide = _blank(prs)
    _header_bar(slide, f"Vânzări pe Agenți & Canale {an}",
                "Val Netă RON")
    _footer(slide)

    items = [(c.get('agent') or c.get('label') or '—', c.get('val_neta') or 0)
             for c in channels_data or []]
    items = [(label, val) for label, val in items if val > 0]
    items.sort(key=lambda x: x[1])

    if items:
        data = CategoryChartData()
        data.categories = [b[0] for b in items]
        data.add_series('Val Netă', [b[1] for b in items])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.3), Inches(1.1),
            Inches(12.7), Inches(5.8),
            data,
        ).chart
        chart.has_legend = False
    return slide


# ── Public API ────────────────────────────────────────────────────────────────

def build_dashboard_ppt(an, cy, py, delta_vn, delta_mb, delta_mn, delta_mpct,
                         agents_data, clients_data, kaufland_pct, bogdan_pct, churn_list,
                         trend_by_year=None, brands_data=None, channels_data=None):
    _check()
    prs = _prs()
    _slide_cover(prs, "Raport Comercial", "Analiză de Performanță", an)
    _slide_kpi_overview(prs, an, cy, py, delta_vn, delta_mb, delta_mn, delta_mpct)
    if trend_by_year:
        _slide_chart_trend(prs, an, trend_by_year)
    if brands_data:
        _slide_chart_brands(prs, an, brands_data)
    if channels_data:
        _slide_chart_channels(prs, an, channels_data)
    _slide_agents_table(prs, an, agents_data)
    _slide_top_clients(prs, an, clients_data)
    _slide_risk(prs, an, kaufland_pct, bogdan_pct, churn_list)
    return _to_bytes(prs)


def build_agent_ppt(agent_name, an, kpi, kpi_py, clients_data, brands_data, skus_data):
    _check()
    prs = _prs()
    _slide_cover(prs, f"Agent: {agent_name}", "Raport Individual", an)
    _slide_agent_detail(prs, agent_name, an, kpi, clients_data, brands_data)

    # Top SKUs slide
    slide = _blank(prs)
    _header_bar(slide, f"{agent_name} — Top Produse {an}")
    _footer(slide)
    s_rows = [{
        "Produs": (r.get('sku') or "—")[:30],
        "Brand": r.get('furnizor') or "—",
        "VN": _fmt_ron(r.get('val_neta')),
        "MB%": _fmt_pct(r.get('marja_bruta_pct')),
        "MN RON": _fmt_ron(r.get('marja_neta')),
        "MN%": _fmt_pct(r.get('marja_neta_pct')),
        "Clienți": str(r.get('nr_clienti') or 0),
    } for r in (skus_data or [])[:15]]
    if s_rows:
        _add_table(slide, s_rows, list(s_rows[0].keys()), 0.3, 1.2, 12.7, 5.8, font_size=8)

    return _to_bytes(prs)


def build_client_ppt(client_name, an, kpi, products_data, yearly_data):
    _check()
    prs = _prs()
    _slide_cover(prs, f"Client: {client_name[:35]}", "Profil de Profitabilitate", an)
    _slide_client_detail(prs, client_name, an, kpi, products_data, yearly_data)
    return _to_bytes(prs)


def build_profitabilitate_ppt(an, agents_data, clients_data, products_data):
    _check()
    prs = _prs()
    _slide_cover(prs, "Analiză Profitabilitate", "Agenți / Clienți / Produse", an)
    _slide_agents_table(prs, an, agents_data)
    _slide_top_clients(prs, an, clients_data, "Top Clienți după Marjă Netă")

    # Products slide
    slide = _blank(prs)
    _header_bar(slide, f"Top Produse după Marjă Netă {an}")
    _footer(slide)
    p_rows = [{
        "#": str(i),
        "Produs": (r.get('sku') or "—")[:28],
        "Brand": r.get('furnizor') or "—",
        "VN": _fmt_ron(r.get('val_neta')),
        "MB%": _fmt_pct(r.get('marja_bruta_pct')),
        "MN RON": _fmt_ron(r.get('marja_neta')),
        "MN%": _fmt_pct(r.get('marja_neta_pct')),
        "Clienți": str(r.get('nr_clienti') or 0),
    } for i, r in enumerate(products_data[:15], 1)]
    if p_rows:
        _add_table(slide, p_rows, list(p_rows[0].keys()), 0.3, 1.2, 12.7, 5.8, font_size=8)

    return _to_bytes(prs)


def _to_bytes(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def send_ppt(buf, filename):
    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def timestamped_filename(base):
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M")
    return f"{base}_{ts}.pptx"


# ── Raportare Basilur ─────────────────────────────────────────────────────────

def build_basilur_ppt(an, period_label, kpi_total, kpi_per_brand,
                      monthly_data, stoc_per_brand, stoc_detail):
    _check()
    BRANDS = ["Basilur", "KingsLeaf", "Tipson"]
    BRAND_COLORS = [
        RGBColor(0x0d, 0x6e, 0xfd),
        RGBColor(0x19, 0x87, 0x54),
        RGBColor(0xfd, 0x7e, 0x14),
    ]
    MONTHS_EN = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    R = 4.55  # RON → USD

    def _usd(v):
        if v is None:
            return '—'
        u = float(v) / R
        if abs(u) >= 1_000_000:
            return f"${u / 1_000_000:.2f}M"
        if abs(u) >= 1_000:
            return f"${u:,.0f}".replace(',', ' ')
        return f"${u:,.0f}"

    prs = _prs()

    # ── Slide 1: Cover + group KPIs ──
    slide = _blank(prs)
    _add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    _add_rect(slide, 0, 0, 0.25, 7.5, RGBColor(0x0d, 0x6e, 0xfd))
    _add_text(slide, "SALES & STOCK REPORT",
              0.5, 1.2, 12, 0.7, font_size=32, bold=True,
              color=C_WHITE, align=PP_ALIGN.LEFT)
    _add_text(slide, "Basilur Group  •  Basilur  •  KingsLeaf  •  Tipson",
              0.5, 2.0, 12, 0.45, font_size=16,
              color=RGBColor(0x9e, 0xc8, 0xff), align=PP_ALIGN.LEFT)
    _add_text(slide, f"Period: {period_label}",
              0.5, 2.55, 12, 0.35, font_size=12, color=C_MGRAY, align=PP_ALIGN.LEFT)
    _add_text(slide, f"All values in USD  |  Rate: 1 USD = {R} RON  |  Torb Logistic SRL  —  Confidential",
              0.5, 7.05, 12, 0.3, font_size=9, color=C_MGRAY, align=PP_ALIGN.LEFT)

    if kpi_total:
        vn   = kpi_total.get("val_neta") or 0
        dlt  = kpi_total.get("delta_vn")
        clnt = kpi_total.get("clienti_activi") or 0
        sku  = kpi_total.get("nr_sku") or 0
        stoc_val = sum((r.get("valoare_achizitie") or 0) for r in (stoc_per_brand or []))
        dlt_str = (f"{'+' if dlt >= 0 else ''}{dlt:.1f}% vs {an - 1}") if dlt is not None else None
        _kpi_card(slide, 0.5,  3.8, 3.1, "Group Net Sales",
                  _usd(vn), dlt_str, (dlt or 0) >= 0)
        _kpi_card(slide, 3.8,  3.8, 2.8, "Active Clients",  str(clnt))
        _kpi_card(slide, 6.8,  3.8, 2.6, "Active SKUs",     str(sku))
        _kpi_card(slide, 9.6,  3.8, 3.2, "Unsold Stock",    _usd(stoc_val))

    # ── Slide 2: KPI per brand ──
    slide = _blank(prs)
    _header_bar(slide, f"KPI by Brand  —  {period_label}",
                f"Basilur Group  |  {an} vs {an - 1}")
    _footer(slide)

    kpi_map = {r["furnizor"]: r for r in (kpi_per_brand or [])}
    stoc_map_s2 = {r["furnizor"]: r for r in (stoc_per_brand or [])}
    card_x = [0.3, 4.6, 8.9]

    for col_i, brand in enumerate(BRANDS):
        k  = kpi_map.get(brand, {})
        st = stoc_map_s2.get(brand, {})
        vn   = k.get("val_neta") or 0
        dlt  = k.get("delta_vn")
        clnt = k.get("clienti_activi") or 0
        nsku = k.get("nr_sku") or 0
        sv   = st.get("valoare_achizitie") or 0
        cx   = card_x[col_i]
        card_w = 4.1
        _add_rect(slide, cx, 1.15, card_w, 0.4, BRAND_COLORS[col_i])
        _add_text(slide, brand.upper(), cx + 0.1, 1.18, card_w - 0.2, 0.35,
                  font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
        dlt_str = (f"{'+' if (dlt or 0) >= 0 else ''}{dlt:.1f}%") if dlt is not None else None
        _kpi_card(slide, cx, 1.65, card_w, "Net Sales",
                  _usd(vn), f"{dlt_str} vs {an - 1}" if dlt_str else None, (dlt or 0) >= 0)
        _kpi_card(slide, cx,              3.5, card_w / 2 - 0.05, "Clients",     str(clnt))
        _kpi_card(slide, cx + card_w / 2, 3.5, card_w / 2 - 0.05, "Active SKUs", str(nsku))
        _kpi_card(slide, cx, 5.15, card_w, "Stock at Cost", _usd(sv))

    # ── Slide 3: Monthly sales evolution ──
    slide = _blank(prs)
    _header_bar(slide, f"Monthly Net Sales Evolution  —  {an}",
                "Net Sales USD per brand per month")
    _footer(slide)

    if monthly_data and any(monthly_data.values()):
        data = CategoryChartData()
        data.categories = MONTHS_EN
        for brand in BRANDS:
            vals = monthly_data.get(brand, [0] * 12)
            data.add_series(brand, tuple(float(v or 0) / R for v in vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.3), Inches(1.15),
            Inches(12.7), Inches(5.8),
            data,
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.plots[0].vary_by_categories = False
        for s_i, series in enumerate(chart.series):
            if s_i < len(BRAND_COLORS):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = BRAND_COLORS[s_i]
    else:
        _add_text(slide, "No sales data available for this period.",
                  0.5, 3.5, 12, 1, font_size=14, color=C_MGRAY)

    # ── Slide 4: Unsold stock at acquisition cost ──
    slide = _blank(prs)
    _header_bar(slide, "Unsold Stock  —  Acquisition Value (USD)",
                "Basilur  •  KingsLeaf  •  Tipson")
    _footer(slide)

    stoc_map = {r["furnizor"]: r for r in (stoc_per_brand or [])}
    total_stoc_val = sum((r.get("valoare_achizitie") or 0) for r in (stoc_per_brand or []))

    for col_i, brand in enumerate(BRANDS):
        st  = stoc_map.get(brand, {})
        val = st.get("valoare_achizitie") or 0
        sku = st.get("nr_sku") or 0
        buc = st.get("total_unitati") or 0
        pct = round(val * 100 / total_stoc_val, 1) if total_stoc_val else 0
        cx  = card_x[col_i]
        _add_rect(slide, cx, 1.15, 4.1, 0.3, BRAND_COLORS[col_i])
        _add_text(slide, brand, cx + 0.08, 1.17, 4.0, 0.25,
                  font_size=11, bold=True, color=C_WHITE)
        _kpi_card(slide, cx, 1.55, 4.1, "Acquisition Value", _usd(val))
        _kpi_card(slide, cx,      3.15, 1.95, "SKU Count",    str(sku))
        _kpi_card(slide, cx + 2.1, 3.15, 1.95, "% of total", f"{pct}%")
        _kpi_card(slide, cx, 4.75, 4.1,  "Units",
                  f"{int(buc):,}".replace(",", " "))

    _add_rect(slide, 0.3, 6.35, 12.7, 0.6, C_ACCENT)
    _add_text(slide, f"TOTAL BASILUR GROUP STOCK: {_usd(total_stoc_val)}",
              0.5, 6.4, 12, 0.5, font_size=14, bold=True, color=C_WHITE,
              align=PP_ALIGN.CENTER)

    # ── Slide 5: Top SKUs by stock value ──
    if stoc_detail:
        slide = _blank(prs)
        _header_bar(slide, "Top SKU Unsold Stock  —  Acquisition Value",
                    f"Top {min(len(stoc_detail), 20)} SKUs by value (USD)")
        _footer(slide)
        top = stoc_detail[:20]
        tbl_rows = [{
            "Brand":       r.get("furnizor") or "?",
            "Code":        str(r.get("cod_produs") or "?"),
            "SKU":         (r.get("sku") or "?")[:35],
            "Qty":         f"{int(r.get('cantitate') or 0):,}".replace(",", " "),
            "Unit Cost":   _usd(r.get("pret_achizitie")),
            "Total Value": _usd(r.get("valoare_achizitie")),
            "Days":        str(r.get("nr_zile_stoc") or "?"),
        } for r in top]
        if tbl_rows:
            _add_table(slide, tbl_rows, list(tbl_rows[0].keys()),
                       0.3, 1.2, 12.7, 5.8, font_size=7.5)

    return _to_bytes(prs)
